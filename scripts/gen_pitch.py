#!/usr/bin/env python3
"""Genera recursos/presentacion.pptx — pitch DatosVivos, Datos al Ecosistema 2026.

12 diapositivas de pitch + 6 técnicas de respaldo. Identidad gov.co/DatosVivos:
azul institucional #004884, blanco, semáforo gov.co. Diseño GRÁFICO: tarjetas,
flujos, barras y cifras grandes — el texto corrido vive en los docs, no aquí.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

AZUL = RGBColor(0x00, 0x48, 0x84)
AZUL2 = RGBColor(0x33, 0x66, 0xCC)
AZUL_SUAVE = RGBColor(0xDC, 0xE6, 0xF2)
INK = RGBColor(0x20, 0x21, 0x24)
MUTED = RGBColor(0x5F, 0x63, 0x68)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CREMA = RGBColor(0xF4, 0xF6, 0xF9)
OK = RGBColor(0x06, 0x84, 0x60)
WARN = RGBColor(0xFF, 0xAB, 0x00)
BAD = RGBColor(0xC3, 0x2D, 0x4B)
CELESTE = RGBColor(0xBF, 0xD4, 0xE8)

W, H = Inches(13.333), Inches(7.5)
FONT = "Nunito Sans"
MONO = "IBM Plex Mono"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
N = [0]


def slide(bg=BLANCO):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    N[0] += 1
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size, color=INK, bold=False, first=None, font=FONT,
         align=PP_ALIGN.LEFT, space_after=6):
    p = tf.paragraphs[0] if first is None and not tf.paragraphs[0].runs else tf.add_paragraph()
    if first is True:
        p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    p.alignment = align
    p.space_after = Pt(space_after)
    return p


def shape(s, kind, x, y, w, h, color, line_color=None):
    sh = s.shapes.add_shape(kind, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line_color is not None:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def rect(s, x, y, w, h, color, line=False):
    return shape(s, MSO_SHAPE.RECTANGLE, x, y, w, h,
                 color, AZUL if line else None)


def card(s, x, y, w, h, fill=CREMA, accent=None):
    """Tarjeta con banda superior de acento opcional. Devuelve el text_frame."""
    r = rect(s, x, y, w, h, fill)
    if accent is not None:
        rect(s, x, y, w, Inches(0.09), accent)
    tf = r.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)
    return tf


def footer(s, dark=False):
    tf = box(s, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    c = CELESTE if dark else MUTED
    para(tf, f"DatosVivos · datosvivos.co  —  Datos al Ecosistema 2026 · Equipo 93 · Reto 7 (id 102) · Nivel Avanzado        {N[0]:02d}",
         9, color=c, font=MONO)


def kicker_title(s, kicker, title, title_size=32, color=INK, kcolor=AZUL2):
    tf = box(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.5))
    para(tf, kicker.upper(), 12, color=kcolor, bold=True, font=MONO, first=True)
    para(tf, title, title_size, color=color, bold=True, space_after=0)


def stat_cards(s, stats, y=Inches(2.1), card_h=Inches(1.9), num_size=30):
    n = len(stats)
    gap = Inches(0.25)
    total_w = W - Inches(1.2)
    cw = Emu(int((total_w - gap * (n - 1)) / n))
    x = Inches(0.6)
    for num, label, color in stats:
        r = rect(s, x, y, cw, card_h, CREMA)
        rect(s, x, y, cw, Inches(0.09), color)
        tf = r.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, num, num_size, color=color, bold=True, font=MONO,
             align=PP_ALIGN.CENTER, first=True, space_after=2)
        para(tf, label, 12, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)
        x = Emu(int(x + cw + gap))


def flow(s, steps, y, h=Inches(1.5), x0=Inches(0.6), x1=None, color=AZUL,
         text_color=BLANCO, size=13):
    """Pipeline horizontal: chevrons de fondo + TEXTBOX superpuesto (el texto
    dentro del chevron colapsa en vertical por el notch del shape)."""
    x1 = x1 or (W - Inches(0.6))
    n = len(steps)
    gap = Inches(0.12)
    cw = Emu(int((x1 - x0 - gap * (n - 1)) / n))
    x = x0
    for i, (titulo, sub) in enumerate(steps):
        kind = MSO_SHAPE.CHEVRON if i > 0 else MSO_SHAPE.PENTAGON
        sh = shape(s, kind, x, y, cw, h, color)
        sh.text_frame.paragraphs[0].add_run().text = ""  # sin texto interno
        # overlay: caja de texto centrada sobre el cuerpo del chevron
        notch = Emu(int(cw * 0.22))
        tb = s.shapes.add_textbox(Emu(int(x + notch if i else x + Inches(0.05))),
                                  y, Emu(int(cw - notch * 2.2)), h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        para(tf, titulo, size, color=text_color, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=1)
        if sub:
            for linea in sub.split("\n"):
                para(tf, linea, size - 3,
                     color=CELESTE if color == AZUL else MUTED,
                     align=PP_ALIGN.CENTER, space_after=0)
        x = Emu(int(x + cw + gap))


def stacked_bar(s, x, y, w, h, segments, label_size=12):
    """Barra apilada horizontal: [(frac, color, etiqueta), ...]."""
    total = sum(f for f, _, _ in segments)
    cx = x
    for frac, color, _ in segments:
        seg_w = Emu(int(w * (frac / total)))
        rect(s, cx, y, seg_w, h, color)
        cx = Emu(int(cx + seg_w))
    # leyenda debajo (solo si hay etiquetas)
    if not any(et for _, _, et in segments):
        return
    tf = box(s, x, Emu(int(y + h + Inches(0.08))), w, Inches(0.4))
    p = tf.paragraphs[0]
    for frac, color, et in segments:
        r = p.add_run()
        r.text = "■ "
        r.font.color.rgb = color
        r.font.size = Pt(label_size + 2)
        r.font.name = MONO
        r2 = p.add_run()
        r2.text = f"{et}      "
        r2.font.color.rgb = INK
        r2.font.size = Pt(label_size)
        r2.font.name = FONT


def grid_cards(s, items, cols, y, card_h, gap_x=Inches(0.25),
               gap_y=Inches(0.22), x0=Inches(0.6), title_size=14,
               body_size=12, accent=AZUL2):
    """Grilla de tarjetas [(titulo, cuerpo)] con acento superior."""
    total_w = W - x0 * 2
    cw = Emu(int((total_w - gap_x * (cols - 1)) / cols))
    for i, (titulo, cuerpo) in enumerate(items):
        r, c = divmod(i, cols)
        x = Emu(int(x0 + c * (cw + gap_x)))
        yy = Emu(int(y + r * (card_h + gap_y)))
        tf = card(s, x, yy, cw, card_h, accent=accent)
        para(tf, titulo, title_size, color=AZUL, bold=True, first=True,
             space_after=3)
        if cuerpo:
            para(tf, cuerpo, body_size, color=INK, space_after=0)


def big_num_row(s, items, y, h=Inches(1.05), x0=Inches(0.6), num_size=26,
                label_size=12):
    """Fila de número grande + etiqueta al lado, sin tarjeta (aire)."""
    total_w = W - x0 * 2
    n = len(items)
    cw = Emu(int(total_w / n))
    x = x0
    for num, label, color in items:
        tf = box(s, x, y, cw, h)
        para(tf, num, num_size, color=color, bold=True, font=MONO, first=True,
             space_after=1)
        para(tf, label, label_size, color=MUTED, space_after=0)
        x = Emu(int(x + cw))


# ============================== PITCH ==============================

# 1 · Portada
s = slide(AZUL)
rect(s, Inches(0), Inches(6.85), W, Inches(0.05), AZUL2)
tf = box(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(3.4))
para(tf, "Datos|Vivos", 64, color=BLANCO, bold=True, font=MONO, first=True, space_after=4)
para(tf, "── el panorama de los datos abiertos de Colombia", 20, color=CELESTE, font=MONO, space_after=20)
para(tf, "Datos del Estado, en tus palabras.", 28, color=BLANCO, bold=True)
# mini-barra semáforo como firma visual
stacked_bar(s, Inches(0.95), Inches(5.35), Inches(6.5), Inches(0.28), [
    (9, OK, ""), (20, WARN, ""), (71, BAD, "")])
tf2 = box(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.0))
para(tf2, "Concurso Datos al Ecosistema 2026: IA para Colombia — MinTIC", 15, color=BLANCO, first=True, space_after=2)
para(tf2, "Equipo 93 · Reto de Innovación y Tecnología (Reto 7, id 102) · Nivel Avanzado · GIT TIC — ANI", 13, color=CELESTE, font=MONO)
footer(s, dark=True)

# 2 · El dolor — 4 tarjetas persona con número
s = slide()
kicker_title(s, "El problema", "Nadie tiene el panorama")
dolores = [
    ("LA ENTIDAD", "no sabe cuántos datasets tiene\nni cuántos están al día"),
    ("EL SECTOR", "sin vista consolidada de sus\nentidades adscritas: cero control"),
    ("EL MINTIC", "los portales federados viven\nseparados — no hay agregado"),
    ("EL CIUDADANO", "necesita APIs y SQL\npara una simple cifra"),
]
x = Inches(0.6)
cw, gap = Inches(3.0), Inches(0.24)
for i, (quien, dolor) in enumerate(dolores):
    tf = card(s, x, Inches(2.15), cw, Inches(3.1), accent=BAD)
    para(tf, "✕", 40, color=BAD, bold=True, font=MONO, first=True,
         align=PP_ALIGN.CENTER, space_after=6)
    para(tf, quien, 15, color=INK, bold=True, font=MONO,
         align=PP_ALIGN.CENTER, space_after=6)
    for linea in dolor.split("\n"):
        para(tf, linea, 13, color=MUTED, align=PP_ALIGN.CENTER, space_after=1)
    x = Emu(int(x + cw + gap))
tf = box(s, Inches(0.65), Inches(5.7), Inches(12), Inches(0.7))
para(tf, "El dato público existe — pero nadie puede verlo completo ni usarlo sin conocimientos técnicos.",
     16, color=INK, bold=True, first=True, align=PP_ALIGN.CENTER)
footer(s)

# 3 · La tesis
s = slide(CREMA)
rect(s, Inches(0.9), Inches(2.2), Inches(0.12), Inches(2.6), AZUL)
tf = box(s, Inches(1.35), Inches(2.3), Inches(11), Inches(2.6))
para(tf, "“Si no conocemos, no podemos medir.", 40, color=AZUL, bold=True, first=True, space_after=2)
para(tf, "Y si no medimos, no podemos mejorar.”", 40, color=AZUL, bold=True, space_after=16)
para(tf, "El dato bien gobernado es infraestructura — tan estratégica como las vías.", 18, color=INK)
footer(s)

# 4 · La solución: 3 niveles con flujo
s = slide()
kicker_title(s, "La solución", "Tres niveles, una sola fuente de verdad")
levels = [
    ("1", "INICIO", "datosvivos.co", "El panorama nacional en vivo.\nPara decidir en segundos."),
    ("2", "DETALLE ENTIDAD", "/tablero · Power BI", "Filtros por sector, entidad\ny territorio. Para gestionar."),
    ("3", "BUSCAR", "/buscar · lenguaje natural", "Pregunta en tus palabras;\ncifra verificada con fuente."),
]
x = Inches(0.6)
cw, gap = Inches(3.85), Inches(0.35)
for num, kick, ruta, desc in levels:
    # círculo numerado
    circ = shape(s, MSO_SHAPE.OVAL, Emu(int(x + cw / 2 - Inches(0.42))),
                 Inches(1.95), Inches(0.84), Inches(0.84), AZUL)
    tfc = circ.text_frame
    tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tfc, num, 30, color=BLANCO, bold=True, font=MONO,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    tf = card(s, x, Inches(2.55), cw, Inches(2.9), accent=AZUL2)
    para(tf, "", 6, first=True, space_after=8)  # aire bajo el círculo
    para(tf, kick, 17, color=AZUL, bold=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf, ruta, 13, color=AZUL2, font=MONO, align=PP_ALIGN.CENTER, space_after=8)
    for linea in desc.split("\n"):
        para(tf, linea, 13, color=INK, align=PP_ALIGN.CENTER, space_after=1)
    x = Emu(int(x + cw + gap))
tf = box(s, Inches(0.65), Inches(5.8), Inches(12), Inches(0.6))
para(tf, "De lo general a lo puntual: panorama → detalle → dato exacto. Cada nivel dirige al siguiente.",
     14, color=MUTED, first=True, align=PP_ALIGN.CENTER)
footer(s)

# 5 · Demo panorama — cifras + semáforo visual
s = slide()
kicker_title(s, "Demo 1/3 · Inicio", "El panorama, en vivo — se actualiza solo cada día")
stat_cards(s, [
    ("25.226", "datasets integrados", AZUL),
    ("1.423", "entidades publicadoras", AZUL),
    ("6", "portales consolidados", AZUL),
    ("29", "variables curadas c/u", AZUL),
], y=Inches(1.95), card_h=Inches(1.55), num_size=32)
tf = box(s, Inches(0.65), Inches(3.9), Inches(12), Inches(0.5))
para(tf, "EL HALLAZGO QUE NADIE VEÍA — frescura contra la promesa de cada entidad:",
     13, color=AZUL2, bold=True, font=MONO, first=True)
stacked_bar(s, Inches(0.65), Inches(4.45), Inches(12.03), Inches(0.75), [
    (9, OK, "9 % al día"),
    (20, WARN, "20 % atrasado"),
    (71, BAD, "71 % MUY atrasado (rojo)"),
], label_size=13)
tf = box(s, Inches(0.65), Inches(5.75), Inches(12), Inches(0.9))
para(tf, "El 71 % del catálogo incumple la frecuencia que su PROPIA entidad declaró — invisible hasta ahora.",
     15, color=INK, bold=True, first=True)
para(tf, "Cifras del corte 2026-07-12, verificables en vivo: datosvivos.co/api/v1/stats/panorama",
     11, color=MUTED, font=MONO)
footer(s)

# 6 · Tablero del decisor — 4 tarjetas
s = slide()
kicker_title(s, "Demo 2/3 · Detalle entidad (Power BI)", "El control de gestión, en dos clics")
grid_cards(s, [
    ("Semáforo por entidad", "verde ≤ frecuencia declarada · amarillo ≤ 2× · rojo > 2×. El % de cumplimiento listo para gestión."),
    ("Uso real", "descargas, vistas e interés reciente: qué datasets le importan a la ciudadanía."),
    ("Territorio DIVIPOLA", "mapa con drill departamento → municipio (~89 % de cobertura inferida con IA)."),
    ("Filtros del decisor", "sector · entidad · acceso · calidad (Ley 1712) · territorio."),
], cols=2, y=Inches(2.0), card_h=Inches(1.75), title_size=16, body_size=13)
tf = box(s, Inches(0.65), Inches(6.0), Inches(12), Inches(0.7))
para(tf, "La IA está DETRÁS del tablero: depuración, consolidación y casos de calidad de los datos se definieron con IA.",
     14, color=AZUL, bold=True, first=True, align=PP_ALIGN.CENTER)
footer(s)

# 7 · El ciudadano pregunta — pipeline visual
s = slide()
kicker_title(s, "Demo 3/3 · Buscar (lenguaje natural)", "“¿Cuántos colegios públicos hay en Boyacá?”")
flow(s, [
    ("Pregunta", "en tus palabras"),
    ("La IA genera", "SQL sobre columnas reales"),
    ("El código verifica", "antes de ejecutar"),
    ("Datos oficiales", "filas reales, no memoria"),
    ("Cifra + fuente", "verificada y citada"),
], y=Inches(2.3), h=Inches(1.6))
big_num_row(s, [
    ("0", "cifras inventadas — si no se puede verificar, se rehúsa", OK),
    ("100 %", "respuestas con fuente citada y consulta visible", AZUL),
    ("Voz", "entrada hablada y respuesta narrada (Ley 1618)", AZUL2),
], y=Inches(4.6), h=Inches(1.2), num_size=30)
footer(s)

# 8 · Dónde vive la IA — diagrama central + dos columnas cortas
s = slide()
kicker_title(s, "Tecnologías emergentes · IA", "La IA razona, el motor verifica")
# banda central con el principio
band = rect(s, Inches(0.6), Inches(1.85), Inches(12.13), Inches(0.85), AZUL)
tfb = band.text_frame
tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tfb, "IA generativa + arquitectura híbrida — pertinente, aplicable e interpretable. Nunca superficial.",
     15, color=BLANCO, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
for x0, titulo, items in [
    (Inches(0.6), "EN LOS TABLEROS", [
        "Consolidación de 6 portales heterogéneos",
        "Clasificación automática de calidad (Ley 1712)",
        "Curación de columnas con LLM + heurísticas",
        "Inferencia territorial DIVIPOLA (~89 %)",
        "Guardas anti-basura en la ingesta"]),
    (Inches(6.85), "EN EL BUSCADOR", [
        "NL2SQL verificado en 3 capas",
        "Embeddings e5 + ChromaDB (retrieval semántico)",
        "6 tipos de respuesta (conteo, suma, mapa…)",
        "Narrativa anti-alucinación (censura números ajenos)",
        "MCP server para cualquier agente de IA"]),
]:
    tf = card(s, x0, Inches(2.95), Inches(5.88), Inches(3.45), accent=AZUL2)
    para(tf, titulo, 14, color=AZUL2, bold=True, font=MONO, first=True, space_after=8)
    for t in items:
        para(tf, "—  " + t, 13.5, color=INK, space_after=7)
footer(s)

# 9 · Metodología — flujo CRISP-ML
s = slide()
kicker_title(s, "Rigor técnico", "CRISP-ML(Q), adaptado a un catálogo vivo")
flow(s, [
    ("Entender", "el problema\ny las fuentes"),
    ("Preparar", "la IA depura\ny consolida"),
    ("Modelar", "generar →\nVERIFICAR"),
    ("Evaluar", "golden sets +\nciclo ciudadano"),
    ("Desplegar", "producción\nreproducible"),
    ("Monitorear", "la Q: deriva del\nCATÁLOGO"),
], y=Inches(2.2), h=Inches(1.7), size=14)
grid_cards(s, [
    ("La adaptación central", "en vez de entrenar-y-congelar un modelo, CADA consulta pasa por generación (LLM) → verificación (código) → ejecución."),
    ("Evaluación doble", "golden sets técnicos versionados + ciclo ciudadano: 50 preguntas reales con respuesta esperada pre-registrada, corridas por ciclos."),
], cols=2, y=Inches(4.55), card_h=Inches(1.6), title_size=15, body_size=13)
footer(s)

# 10 · Nivel Avanzado — grilla 3×2 con la letra del TDR
s = slide()
kicker_title(s, "Cumplimiento del TDR", "Nivel Avanzado, con la letra del pliego")
grid_cards(s, [
    ("Agentes de IA", "consulta y procesa datos abiertos automáticamente para responder a la ciudadanía."),
    ("IA generativa conversacional", "buscador NL2SQL / Text-to-SQL generativo verificado."),
    ("Arquitectura híbrida", "LLM razona + motor determinista verifica + embeddings neuronales."),
    ("Grandes volúmenes", "25.226 datasets · 6 fuentes · 3 protocolos (el intermedio pide 3-10 conjuntos)."),
    ("Estructurados y no estructurados", "metadata + texto libre con embeddings; 29 variables curadas (el intermedio pide 10-20)."),
    ("Despliegue funcional", "actualización diaria automática y EN PRODUCCIÓN: datosvivos.co."),
], cols=3, y=Inches(2.0), card_h=Inches(1.95), title_size=14.5, body_size=12,
    accent=OK)
footer(s)

# 11 · Impacto — cadena de valor + equipo
s = slide()
kicker_title(s, "Impacto y escalabilidad", "De la entidad al país")
flow(s, [
    ("Entidad", "ve su rezago\n→ corrige"),
    ("Sector", "control consolidado\n→ gestiona"),
    ("MinTIC", "panorama medible\n→ política"),
    ("Ciudadanía", "dato sin barreras\n→ apropiación"),
], y=Inches(2.15), h=Inches(1.6), size=15)
tf = box(s, Inches(0.65), Inches(4.05), Inches(12), Inches(0.65))
para(tf, "Escala: agregar un portal es configuración — el patrón sirve a cualquier país con Socrata, CKAN o DCAT.",
     14, color=INK, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = card(s, Inches(0.6), Inches(4.95), Inches(12.13), Inches(1.6), accent=AZUL)
para(tf, "EQUIPO GIT TIC — AGENCIA NACIONAL DE INFRAESTRUCTURA (ANI)", 13,
     color=AZUL2, bold=True, font=MONO, first=True, space_after=6)
para(tf, "Hernán Darío Gutiérrez Casas — líder estratégico  ·  Ileana Andrea Navarro Castrillón — líder de equipo y comunicaciones  ·  Jhonatan Sneider Rico Pinto — líder técnico y de datos",
     14, color=INK, space_after=0)
footer(s)

# 12 · Cierre + roadmap
s = slide(AZUL)
tf = box(s, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.6))
para(tf, "LO QUE YA CORRE · JULIO 2026", 13, color=CELESTE, bold=True, font=MONO, first=True)
hechos = [
    ("~1.5 s", "motor de lenguaje\nen la API de Claude"),
    ("10.279", "datasets en bodega local\n→ respuesta en milisegundos"),
    ("100 %", "catalogación\ntemática"),
    ("50", "preguntas ciudadanas\n→ mejoras estructurales"),
]
x = Inches(0.9)
cw, gap = Inches(2.85), Inches(0.25)
for num, label in hechos:
    tf = box(s, x, Inches(1.95), cw, Inches(1.6))
    para(tf, num, 30, color=BLANCO, bold=True, font=MONO, first=True, space_after=3)
    for linea in label.split("\n"):
        para(tf, linea, 12, color=CELESTE, space_after=1)
    x = Emu(int(x + cw + gap))
tf = box(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.8))
para(tf, "Lo que sigue: filtros de año/municipio dentro del dataset · respuestas compuestas (KPI + tendencia + per cápita).",
     14, color=BLANCO, first=True)
rect(s, Inches(0.9), Inches(4.75), Inches(11.5), Inches(0.03), AZUL2)
tf = box(s, Inches(0.9), Inches(5.05), Inches(11.5), Inches(1.6))
para(tf, "datosvivos.co", 40, color=BLANCO, bold=True, font=MONO, first=True, space_after=4)
para(tf, "El panorama de los datos abiertos de Colombia. En vivo, verificable, para decidir.", 16, color=CELESTE)
footer(s, dark=True)

# ========================= RESPALDO TÉCNICO =========================

def backup(kicker, title):
    s = slide()
    kicker_title(s, "Respaldo técnico · " + kicker, title, title_size=28)
    return s


def bullets(s, items, x=Inches(0.65), y=Inches(2.0), w=Inches(12.0),
            h=Inches(4.7), size=17, gap=10):
    tf = box(s, x, y, w, h)
    for i, it in enumerate(items):
        head, body = it
        para(tf, head, size, color=AZUL, bold=True, first=(i == 0), space_after=2)
        para(tf, body, size - 2, color=INK, space_after=gap)


# B1 arquitectura — capas como flujo vertical
s = backup("B1", "Arquitectura de 3 capas (fiel al plan de inscripción)")
capas = [
    ("PRESENTACIÓN", "Next.js (panorama + buscador SSE) · Power BI publish-to-web sobre CSVs públicos de la API", AZUL2),
    ("MOTOR", "FastAPI + PostgreSQL (catálogo curado, vistas _decisor) · ChromaDB (retrieval semántico en ambos caminos) · DuckDB (bodega Parquet 10.279 datasets + CSVs federados) · LLM intercambiable (prod: API de Claude)", AZUL),
    ("MCP / FUENTES", "MCP server sobre las APIs de datos.gov.co (Discovery · SODA · Metadata) + cosecha CKAN/DCAT — search · metadata · query · cross", INK),
]
y = Inches(2.0)
for titulo, desc, color in capas:
    r = rect(s, Inches(0.6), y, Inches(12.13), Inches(1.35), CREMA)
    rect(s, Inches(0.6), y, Inches(0.09), Inches(1.35), color)
    tf = r.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, titulo, 14, color=color, bold=True, font=MONO, first=True, space_after=3)
    para(tf, desc, 12.5, color=INK, space_after=0)
    y = Emu(int(y + Inches(1.35) + Inches(0.18)))
tf = box(s, Inches(0.65), Inches(6.55), Inches(12), Inches(0.5))
para(tf, "Despliegue: Docker Compose reproducible · túnel seguro de salida · cron diario (ETL + regla de cola de la bodega).",
     12, color=MUTED, first=True, font=MONO)
footer(s)

# B2 motor verificado — flujo de 3 capas
s = backup("B2", "Motor NL2SQL: generar no basta — hay que verificar")
flow(s, [
    ("Capa 1", "el LLM genera viendo\nSOLO columnas reales"),
    ("Capa 2", "código valida columnas,\nfunciones y solo-lectura"),
    ("Capa 3", "cada cifra se contrasta\ncontra las filas reales"),
    ("Repara / Rehúsa", "si no verifica,\nNO responde"),
], y=Inches(2.2), h=Inches(1.7), size=14)
grid_cards(s, [
    ("Preferimos no responder a responder mal", "ciclo de reparación acotado; si la verificación persiste en fallar, refusal explícito."),
    ("Respaldo determinista", "plantillas por tipo (conteo · suma · comparación · ranking · tendencia · mapa) como camino estructurado con chips."),
], cols=2, y=Inches(4.55), card_h=Inches(1.6), title_size=15, body_size=13)
footer(s)

# B3 modelo de datos
s = backup("B3", "Modelo de datos: una vista curada como fuente única")
bullets(s, [
    ("datasets (42 columnas)", "una fila por dataset del catálogo integrado; upsert idempotente por dataset_id (auditado: 0 duplicados de clave)."),
    ("v_dataset_status_decisor (29 variables)", "identidad, semáforo de frescura, uso, acceso, territorio DIVIPOLA y calidad."),
    ("v_entity_summary_decisor", "agregado por entidad con pct_verdes: el indicador de cumplimiento listo para gestión."),
    ("dataset_snapshots (bodega)", "manifest de los 10.279 Parquet locales: frescura del snapshot decide bodega-vs-vivo; orígenes muertos (403) penalizan el ranking."),
    ("Una sola fuente de verdad", "panorama web, tablero Power BI y CSVs públicos leen LA MISMA vista — imposible que se desalineen."),
])
footer(s)

# B4 calidad
s = backup("B4", "Calidad de datos: medida, no declarada")
bullets(s, [
    ("Auditoría columna a columna", "17/18 columnas al 100 % de fidelidad contra la fuente Socrata (reportes versionados en eval/reports/)."),
    ("Clasificación continua", "los reportes administrativos (Ley 1712) se separan de los datos temáticos automáticamente, en cada corrida del ETL."),
    ("Catalogación temática al 100 %", "0 datasets útiles sin categoría; vocabulario consolidado a 26 canónicas, re-unificado cada noche."),
    ("Inferencia territorial", "códigos DIVIPOLA asignados por IA con confianza registrada (~89 % de cobertura)."),
    ("Guardas anti-basura", "metadata de plantilla sin diligenciar ({{...}}) se limpia en TODOS los campos de la ingesta."),
])
footer(s)

# B5 evaluación
s = backup("B5", "Evaluación y trazabilidad")
bullets(s, [
    ("Golden sets versionados", "eval/golden_queries.yaml y golden_chips.yaml — corridas reproducibles (últimas: 18/18 y 38/38, 0 falsos verificados)."),
    ("Ciclo ciudadano", "50 preguntas reales con respuesta esperada PRE-registrada (eval/ciudadano/): los patrones detectados se vuelven mejoras estructurales, y el ciclo se repite."),
    ("36 archivos de pruebas", "verificador SoQL, validador anti-alucinación, reparación, geo/DIVIPOLA, cosecha, MCP, rutas de API."),
    ("Decisiones documentadas", "ADRs públicos (017-023): del principio 'la IA razona, el motor verifica' al pivote panorama-decisor."),
    ("Telemetría anónima", "lo más consultado alimenta la mejora sin datos personales."),
])
footer(s)

# B6 soberanía
s = backup("B6", "Soberanía, seguridad y cumplimiento")
bullets(s, [
    ("Infraestructura del Estado", "el servicio corre en infraestructura pública; los datos consultados son públicos por definición."),
    ("Sin datos personales", "cero registro, cero rastreadores; telemetría agregada y anónima."),
    ("Solo lectura", "el motor no puede escribir en ninguna fuente; verificación de solo-lectura en la capa 2."),
    ("Repo abierto y auditable", "código, documentación, pruebas y evaluación públicos — replicable con Docker (guía de validación para el jurado)."),
    ("Accesibilidad", "Ley 1618 de 2013 y WCAG 2.1 AA: voz, narración, contraste, escala tipográfica."),
])
footer(s)

OUT = "recursos/presentacion.pptx"
prs.save(OUT)
print(f"OK → {OUT} ({len(prs.slides._sldIdLst)} diapositivas)")
